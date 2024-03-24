import re, os
import requests
from bs4 import BeautifulSoup
from urllib.parse import urljoin
import ebooklib
from ebooklib import epub
import warnings
import whisper
import shutil
import google.generativeai as genai
from PyPDF2 import PdfReader
from odf import text, teletype
import docx
from striprtf.striprtf import rtf_to_text
from striprtf.striprtf import rtf_to_text
from odf.opendocument import load
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
import os
from langchain_google_genai import GoogleGenerativeAIEmbeddings
import google.generativeai as genai
from langchain.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from dotenv import load_dotenv

os.environ["KMP_DUPLICATE_LIB_OK"]="TRUE"

load_dotenv()
os.getenv("GOOGLE_API_KEY")
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))
visited = set()
BASE_DIR = './data'
main = ''
class ConvertToText:
    def __init__(self):
        pass
    def create_directory(self, url):
        try:
            response = requests.get(url)
            response.raise_for_status()

            soup = BeautifulSoup(response.content, 'html.parser')
            elements = soup.find(attrs={'id': 'breadcrumbs'})
            text_content = ""
            if elements:
                for element in elements:
                    text_content += element.get_text(separator='', strip=True)
                dirList = text_content.split('/')
                directory_name = os.path.join(BASE_DIR, '/'.join(dirList[:-1])) + '/'
                file_name = dirList[-1]
                os.makedirs(directory_name, exist_ok=True)  # Create directory if it doesn't exist
            else:
                directory_name = BASE_DIR
                file_name = url.split('/')[-1] 
            return directory_name, file_name
        except requests.exceptions.RequestException as e:
            print("Error creating directory: {e}")
            pass

    def extract_and_save_article_text(self, url):
        """Extracts text from within the article tag, using h1 for filename."""
        try:
            response = requests.get(url)
            response.raise_for_status()

            soup = BeautifulSoup(response.content, 'html.parser')

            # Find the article tag
            article = soup.find('article')
            print(f"Extracting article text from {url}")

            if article:
                # Extract text within the article
                h1_tag = article.find('h1')
                # find tag with id subtitle-bar and remove it
                subtitle_bar = article.find('div', {'id': 'subtitle-bar'})
                breadcrumbs = article.find('div', {'id': 'breadcrumbs'})
                if breadcrumbs:
                    breadcrumbs.decompose()

                if subtitle_bar:
                    subtitle_bar.decompose()

                t_filename = h1_tag.text.strip() if h1_tag else url.split('/')[-1] + '.txt'
                for li in article.find_all('li'):
                    li.replace_with(li.text + ' ')  # Replace the <li> tag with its text and a space
                for note in article.find_all('blockquote'):
                    note.replace_with(note.text + ' ')

                text_content = article.get_text(separator='\n', strip=True)
                print(f"Article text extracted and saved from {url}")
                return text_content, t_filename

            else:
                # remove nav tag from the soup
                nav = soup.find('nav')
                if nav:
                    nav.decompose()
                main = soup.find('main')
                if main:
                    h1_tag = main.find('h1')
                    t_filename = h1_tag.text.strip() if h1_tag else url.split('/')[-1] + '.txt'
                    for li in main.find_all('li'):
                        li.replace_with(li.text + ' ')
                    for note in main.find_all('blockquote'):
                        note.replace_with(note.text + ' ')
                    text_content = main.get_text(separator='\n', strip=True)
                    print(f"Text extracted and saved from {url}")
                    return text_content, t_filename
                body = soup.find('body')
                if body:
                    h1_tag = body.find('h1')
                    t_filename = h1_tag.text.strip() if h1_tag else url.split('/')[-1] + '.txt'
                    for li in body.find_all('li'):
                        li.replace_with(li.text + ' ')
                    for note in body.find_all('blockquote'):
                        note.replace_with(note.text + ' ')
                    text_content = body.get_text(separator='\n', strip=True)
                    print(f"Text extracted and saved from {url}")
                    return text_content, t_filename

        except requests.exceptions.RequestException as e:
            # print(f"Error processing {url}: {e}")
            pass

    def find_links(self, url, base_url, main_url='', depth=None):
        global main
        if depth is not None and depth <= 0:
            return
        if main_url != '':
            main = main_url.split('.')[0]
        if url in visited:
            print(f"{url} is already Visited \n")
            return
        elif main not in url:
            print(f"{url} is not in the main domain {main_url} \n")
            return

        visited.add(url)

        try:
            response = requests.get(url)
            response.raise_for_status()

            soup = BeautifulSoup(response.content, 'html.parser')
            print(f"Processing {url}")
            links = soup.find_all('a', href=True)
            print(f"Found {len(links)} links on {url}")

            text_content, t_filename = self.extract_and_save_article_text(url)  or (None, None)
            if text_content:
                print(f"Creating Directory for {url}")
                directory_name, file_name = self.create_directory(url)
                if directory_name == '':
                    directory_name = 'extracted_articles/'
                    file_name = t_filename
                file_path = os.path.join(directory_name, file_name)
                if os.path.exists(file_path):
                    print(f"File {file_path}.txt already exists, skipping... \n")
                else:
                    print(f"Saving {file_path}")
                    os.makedirs(directory_name, exist_ok=True)
                    text_content = f"Reference Url: {url} \n{text_content}"
                    with open(f'{file_path}.txt', 'w', encoding='utf-8') as f:
                        f.write(text_content)

            for link in links:
                absolute_url = urljoin(base_url, link['href'])  # Construct absolute URLs

                # Filtering logic
                if any(x in absolute_url for x in ["support", "contact-us", "release-info", "search", 'release-notes', 'terms-of-use', 'contribute', 'termsofuse']):
                    continue 
                self.find_links(absolute_url, base_url,depth=depth-1 if depth is not None else None)  

        except requests.exceptions.RequestException as e:
            # print(f"Error processing {url}: {e}")
            pass

    def get_pdf_text(self, pdf):
        text=""
        pdf_reader= PdfReader(pdf)
        for page in pdf_reader.pages:
            text+= page.extract_text()
        return text

    def read_rtf_file(self, file_path):
        with open(file_path) as infile:
            content = infile.read()
            text = rtf_to_text(content)
        return text

    def read_odt_file(self, file_path):
        full_text = []
        odt_doc = load(file_path)
        paragraphs = odt_doc.getElementsByType(text.P)
        for i in paragraphs:
            paragraph = teletype.extractText(i)
            full_text.append(paragraph)
        return '\n'.join(full_text)

    def read_ppt_file(self, file_path):
        prs = Presentation(file_path)
        full_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        # Split text into words and join with a single space to remove extra whitespace
                        text = ' '.join(paragraph.text.split())
                        # Only append non-empty text
                        if text.strip():
                            full_text.append(text)
        # Join paragraphs with a single newline
        return '\n'.join(full_text)

    def read_epub_file(self, file_path):
        # Filter out the ebooklib warning
        with warnings.catch_warnings():
            warnings.filterwarnings("ignore", category=UserWarning)
            book = epub.read_epub(file_path)
        full_text = []
        for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
            soup = BeautifulSoup(item.get_content(), 'html.parser')
            text = soup.get_text()
            # Split text into words and join with a single space to remove extra whitespace
            cleaned_text = ' '.join(text.split())
            if cleaned_text.strip():
                full_text.append(cleaned_text)
        return '\n'.join(full_text)

    def read_py_file(self, file_path):
        with open(file_path, "r", encoding="utf-8") as file:
            return file.read()
        
    def read_audio_file(self, file_path):
        model = whisper.load_model("small")
        result = model.transcribe(file_path)
        return(result["text"])

    def read_word_file(self, file_path):
        doc = docx.Document(file_path)
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        return '\n'.join(full_text)

    def get_link_text(self):
        text = ''
        # get the text from base dir
        for root, dirs, files in os.walk('./data'):
            for name in files:
                with open(os.path.join(root, name), 'r', encoding='utf-8') as file:
                    text += (file.read())
        # delete all the files and folders in the base dir
        if os.path.exists('./data'):
            shutil.rmtree('./data')
        
        return text

    def get_text(self, files):
        for file in files:
            text = ''
            if file.name.endswith(".pdf"):
                text = (self.get_pdf_text(file))
            elif file.name.endswith(".docx"):
                text = (self.read_word_file(file))
            elif file.name.endswith(".rtf"):
                text = (self.read_rtf_file(file))
            elif file.name.endswith(".odt"):
                text = (self.read_odt_file(file))
            elif file.name.endswith(".pptx"):
                text = (self.read_ppt_file(file))
            elif file.name.endswith(".epub"):
                text = (self.read_epub_file(file))
            elif file.name.endswith(".py"):
                text = (self.read_py_file(file))
            elif file.name.endswith(".wav"):
                text = (self.read_audio_file(file))

            # create a text file and store the text in text file
            # what ever the file name is, store the text in the file with the same name but as .txt
            # filename = file.name.split('.')[0]
            # with open(f'./data/{filename}.txt', 'w') as f:
            #     f.write(text)
        return text
    # def main():
    #     starting_url = "https://fnbdocs.apteancloud.com/bc/SUL/contents/"  
    #     find_links(starting_url, starting_url)

    def get_text_chunks(self, text):
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
        chunks = text_splitter.split_text(text)
        return chunks


    def get_vector_store(self, text_chunks, userName):
        # shutil.rmtree("faiss_index", ignore_errors=True)
        embeddings = GoogleGenerativeAIEmbeddings(model = "models/embedding-001")
        vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)


        if os.path.exists(f"./db/{userName}/faiss_index"):
            local_index=FAISS.load_local(f"./db/{userName}/faiss_index", embeddings)
            local_index.merge_from(vector_store)
            local_index.save_local(f"./db/{userName}/faiss_index")
        else:
            vector_store.save_local(f"./db/{userName}/faiss_index")

    def delete_user_data(self, userName):
        shutil.rmtree(f"./db/{userName}", ignore_errors=True)

    def get_conversational_chain(self):

        prompt_template = """
        Answer the question as detailed as possible from the provided context, try to provide the complete content for the query , make sure to provide all the details step by step, if the answer is not in
        provided context just say, "answer is not available in the context", don't provide the wrong answer\n\n
        Context:\n {context}?\n
        Question: \n{question}\n

        Answer:
        """

        model = ChatGoogleGenerativeAI(model="gemini-pro",
                                temperature=0.3)

        prompt = PromptTemplate(template = prompt_template, input_variables = ["context", "question"])
        chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)

        return chain



    def user_input(self, user_question, userName):
        try:
            embeddings = GoogleGenerativeAIEmbeddings(model = "models/embedding-001")
            
            new_db = FAISS.load_local(f"./db/{userName}/faiss_index", embeddings)
            docs = new_db.similarity_search(user_question)

            chain = self.get_conversational_chain()

            
            response = chain(
                {"input_documents":docs, "question": user_question}
                , return_only_outputs=True)

            return response["output_text"]
        except Exception as e:
            print(f"An error occurred: {e if e else 'Provide your data to get the response'}")
            return None

    def get_gemini_response(self, question):
        model=genai.GenerativeModel("gemini-pro") 
        chat = model.start_chat(history=[])
        
        response=chat.send_message(question,stream=True)
        return response